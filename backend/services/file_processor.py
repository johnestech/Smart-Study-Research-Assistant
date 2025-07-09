import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
import io
import mimetypes
from typing import Dict, Any, Optional
from werkzeug.utils import secure_filename
from config import Config

class FileProcessor:
    def __init__(self):
        self.allowed_extensions = Config.ALLOWED_EXTENSIONS
        self.max_file_size = Config.MAX_CONTENT_LENGTH
    
    def is_allowed_file(self, filename: str) -> bool:
        """Check if file extension is allowed"""
        if '.' not in filename:
            return False
        extension = filename.rsplit('.', 1)[1].lower()
        return extension in self.allowed_extensions
    
    def get_file_type(self, filename: str) -> str:
        """Get file type from filename"""
        if '.' not in filename:
            return 'unknown'
        return filename.rsplit('.', 1)[1].lower()
    
    def validate_file(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Validate file size and type"""
        try:
            # Check file size
            if len(file_data) > self.max_file_size:
                return {
                    'valid': False,
                    'error': f'File size exceeds maximum limit of {self.max_file_size // (1024*1024)}MB'
                }
            
            # Check file extension
            if not self.is_allowed_file(filename):
                return {
                    'valid': False,
                    'error': 'File type not supported. Supported types: PDF, Word (.doc, .docx), PowerPoint (.ppt, .pptx), Images (.jpg, .jpeg, .png, .gif)'
                }
            
            return {'valid': True}
        
        except Exception as e:
            return {'valid': False, 'error': f'File validation error: {str(e)}'}
    
    def extract_text_from_pdf(self, file_data: bytes) -> str:
        """Extract text from PDF file"""
        try:
            doc = fitz.open(stream=file_data, filetype="pdf")
            text = ""
            for page_num in range(doc.page_count):
                page = doc[page_num]
                text += page.get_text()
            doc.close()
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting text from PDF: {str(e)}")
    
    def extract_text_from_docx(self, file_data: bytes) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(io.BytesIO(file_data))
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting text from DOCX: {str(e)}")
    
    def extract_text_from_doc(self, file_data: bytes) -> str:
        """Extract text from DOC file (basic support)"""
        try:
            # Note: Full DOC support requires additional libraries like python-docx2txt
            # This is a basic implementation
            return "DOC file processing requires additional setup. Please convert to DOCX format for better support."
        except Exception as e:
            raise Exception(f"Error extracting text from DOC: {str(e)}")
    
    def extract_text_from_pptx(self, file_data: bytes) -> str:
        """Extract text from PPTX file"""
        try:
            prs = Presentation(io.BytesIO(file_data))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting text from PPTX: {str(e)}")
    
    def extract_text_from_ppt(self, file_data: bytes) -> str:
        """Extract text from PPT file (basic support)"""
        try:
            # Note: Full PPT support requires additional libraries
            return "PPT file processing requires additional setup. Please convert to PPTX format for better support."
        except Exception as e:
            raise Exception(f"Error extracting text from PPT: {str(e)}")
    
    def process_image(self, file_data: bytes, filename: str) -> str:
        """Process image file (basic info extraction)"""
        try:
            image = Image.open(io.BytesIO(file_data))
            
            # Basic image information
            info = f"Image file: {filename}\\n"
            info += f"Format: {image.format}\\n"
            info += f"Size: {image.size[0]}x{image.size[1]} pixels\\n"
            info += f"Mode: {image.mode}\\n"
            
            # Check if image has text (OCR would require additional libraries like pytesseract)
            info += "\\nNote: Text extraction from images requires OCR setup. Please ensure any important text content is also provided in document form."
            
            return info
        except Exception as e:
            raise Exception(f"Error processing image: {str(e)}")
    
    def extract_content(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Extract content from file based on file type"""
        try:
            file_type = self.get_file_type(filename)
            content = ""
            
            if file_type == 'pdf':
                content = self.extract_text_from_pdf(file_data)
            elif file_type == 'docx':
                content = self.extract_text_from_docx(file_data)
            elif file_type == 'doc':
                content = self.extract_text_from_doc(file_data)
            elif file_type == 'pptx':
                content = self.extract_text_from_pptx(file_data)
            elif file_type == 'ppt':
                content = self.extract_text_from_ppt(file_data)
            elif file_type in ['jpg', 'jpeg', 'png', 'gif']:
                content = self.process_image(file_data, filename)
            else:
                raise Exception(f"Unsupported file type: {file_type}")
            
            if not content or content.strip() == "":
                content = f"No extractable content found in {filename}. The file may be empty, corrupted, or contain only images/graphics."
            
            return {
                'success': True,
                'content': content,
                'file_type': file_type,
                'filename': filename
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'file_type': self.get_file_type(filename),
                'filename': filename
            }
    
    def get_file_info(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Get basic file information"""
        try:
            file_type = self.get_file_type(filename)
            file_size = len(file_data)
            
            # Get MIME type
            mime_type, _ = mimetypes.guess_type(filename)
            
            return {
                'filename': filename,
                'file_type': file_type,
                'file_size': file_size,
                'mime_type': mime_type,
                'size_mb': round(file_size / (1024 * 1024), 2)
            }
        except Exception as e:
            return {
                'filename': filename,
                'error': str(e)
            }
    
    def create_file_preview(self, content: str, max_length: int = 200) -> str:
        """Create a preview of file content"""
        if len(content) <= max_length:
            return content
        
        # Try to cut at a sentence or word boundary
        truncated = content[:max_length]
        last_sentence = truncated.rfind('.')
        last_space = truncated.rfind(' ')
        
        if last_sentence > max_length - 50:  # If sentence end is close to limit
            return content[:last_sentence + 1] + "..."
        elif last_space > max_length - 20:   # If word boundary is close to limit
            return content[:last_space] + "..."
        else:
            return truncated + "..."

# Global instance
file_processor = FileProcessor()
